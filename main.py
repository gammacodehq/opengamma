from model import invoke_func
import os
import logging
import argparse
from datasets import load_dataset

parser = argparse.ArgumentParser()
parser.add_argument("--log", action="store_true", help="Enable logging")
args = parser.parse_args()

if args.log:
    logging.basicConfig(
        level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
    )
log = logging.getLogger(__name__)


model = "openai/gpt-oss-20b:free"
model = "kwaipilot/kat-coder-pro:free"

dataset = load_dataset("mikeoxmaul/opengamma-prs", streaming=True)
first = dataset["train"].take(1)
task = list(first)[0]["text"]
log.info("Loaded task from dataset")

system_prompt = f"""
Write a complete Python script using the python-pptx library to create a PowerPoint presentation (PPTX).
The script must:
- Use appropriate slide layouts (e.g., title slide, bullet slide).
- Save the presentation to a file named 'test.pptx'.
Output ONLY the Python code as plain text, without markdown, code block markers (e.g., ` + "```" + `python)

Example:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')"""

if __name__ == "__main__":
    result, _ = invoke_func(model, system_prompt, task)
    if result == 1:
        print("Success")
    else:
        print("Failed")
