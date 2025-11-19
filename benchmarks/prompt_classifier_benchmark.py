import os
import sys

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import json
import numpy as np
import logging
import concurrent.futures
import time
from datetime import datetime
from pathlib import Path
from functools import cache
import requests
from dotenv import load_dotenv
from sklearn.linear_model import LogisticRegression
from datasets import load_dataset
from model import invoke_func

load_dotenv()

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
log = logging.getLogger(__name__)

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

system_prompt = f"""
Write a complete Python script using the python-pptx library to create a PowerPoint presentation (PPTX).
The script must:
- Use appropriate slide layouts (e.g., title slide, bullet slide).
- Save the presentation to a file named 'test.pptx'.
Output ONLY the Python code as plain text, without markdown, code block markers (e.g., ` + "```" + `python)

Example:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')"""

# Базовый промпт
basic_prompt = f"""
Generate Python code for python-pptx presentation. OUTPUT MUST BE PURE PYTHON CODE ONLY.

CRITICAL:
- NO comments (# comment)
- NO explanations
- NO markdown (```python or ```)
- NO text outside code
- NO example usage
- NO instructions

REQUIREMENTS:
- Use python-pptx
- Create multiple slides with different layouts  
- Save as 'test.pptx'
- Code must run without errors

NOW GENERATE PURE PYTHON CODE:"""

# Детальный промпт с примерами из документации
detailed_prompt = f"""
Generate Python code for python-pptx presentation. OUTPUT MUST BE PURE PYTHON CODE ONLY.

PYTHON-PPTX SPECIFICS:
- Use Presentation() to create presentation
- Use slide_layouts[] for slide types: 0=Title, 1=Title and Content, 5=Title only
- Use Inches() for measurements
- Add slides with prs.slides.add_slide(layout)
- Set text with shape.text = "text"
- Save with prs.save('test.pptx')

CODE STRUCTURE:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Title"
prs.save('test.pptx')

OUTPUT MUST BE PURE PYTHON CODE ONLY. NO COMMENTS. NO EXPLANATIONS."""

# Минималистичный промпт
minimal_prompt = f"""
Generate python-pptx code. Output ONLY Python code.

from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Title"
prs.save('test.pptx')

Continue with more slides:"""

# Строгий промпт
structured_prompt = f"""
PYTHON CODE ONLY. NO COMMENTS. FOLLOW TEMPLATE:

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "TITLE"
slide.placeholders[1].text = "SUBTITLE"

# ADD MORE SLIDES WITH DIFFERENT LAYOUTS

prs.save('test.pptx')"""

prompts = {
    0: system_prompt,
    1: basic_prompt,
    2: detailed_prompt,
    3: structured_prompt,
}


@cache
def get_embedding(text):
    key = os.getenv("OPENROUTER_KEY")
    response = requests.post(
        "https://openrouter.ai/api/v1/embeddings",
        headers={
            "Authorization": f"Bearer {key}",
            "Content-Type": "application/json",
        },
        json={
            "model": "sentence-transformers/paraphrase-minilm-l6-v2",
            "input": text,
        },
    )
    data = response.json()
    embedding = data["data"][0]["embedding"]
    return embedding


class PromptClassifierBenchmark:
    def __init__(self, model):
        self.model = model
        self.dataset = load_dataset("mikeoxmaul/opengamma-prs-dedup", streaming=True)
        self.tasks = list(self.dataset["train"].skip(100).take(100))
        self.results_dir = Path("results/prompt_classifier")
        self.results_dir.mkdir(parents=True, exist_ok=True)

        # Load classifier
        with open("classifier_tensors.json", "r") as f:
            data = json.load(f)
        self.clf = LogisticRegression()
        self.clf.coef_ = np.array(data["weights"])
        self.clf.intercept_ = np.array(data["bias"])
        self.clf.classes_ = np.array([0, 1, 2, 3])

    def run_task(self, prompt_content, task, i):
        id = 1
        try:
            start_time = time.time()
            result, token_stats = invoke_func(
                self.model, prompt_content, task["text"], id
            )
            execution_time = time.time() - start_time
            return {
                "success": result == 1,
                "time": execution_time,
                "tokens": token_stats,
                "index": i,
            }
        except Exception as e:
            return {
                "success": False,
                "time": 0,
                "tokens": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "error": str(e),
            }

    def benchmark_modes(self):
        results = {}

        # Mode 1: Always use original prompt
        # mode_name = "original"
        # log.info(f"Testing mode: {mode_name}")
        # prompt_content = prompts[0]
        # mode_results = self.run_mode(prompt_content, mode_name)
        # results[mode_name] = mode_results

        # Mode 2: Route using classifier
        mode_name = "routed"
        log.info(f"Testing mode: {mode_name}")
        mode_results = self.run_routed_mode()
        results[mode_name] = mode_results

        return results

    def run_mode(self, prompt_content, mode_name):
        mode_results = {
            "success_count": 0,
            "total_tasks": 0,
            "total_time": 0,
            "token_usage": {
                "prompt_tokens": 0,
                "completion_tokens": 0,
                "total_tokens": 0,
            },
            "errors": [],
            "successful_indices": [],
        }

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            futures = [
                executor.submit(self.run_task, prompt_content, task, i)
                for i, task in enumerate(self.tasks)
            ]
            for future in concurrent.futures.as_completed(futures):
                res = future.result()
                mode_results["total_tasks"] += 1
                mode_results["total_time"] += res["time"]
                mode_results["token_usage"]["prompt_tokens"] += res["tokens"][
                    "prompt_tokens"
                ]
                mode_results["token_usage"]["completion_tokens"] += res["tokens"][
                    "completion_tokens"
                ]
                mode_results["token_usage"]["total_tokens"] += res["tokens"][
                    "total_tokens"
                ]

                log.info(f"Task {res['index'] + 1}/100 complete")
                if res["success"]:
                    mode_results["success_count"] += 1
                    mode_results["successful_indices"].append(res["index"])
                    log.info(f"Success ({res['time']:.2f}s)")
                else:
                    log.info(f"Error ({res['time']:.2f}s)")

                if "error" in res:
                    mode_results["errors"].append(res["error"])
                    log.error(f"Exception in task: {res['error']}")

        if mode_results["total_tasks"] > 0:
            mode_results["success_rate"] = (
                mode_results["success_count"] / mode_results["total_tasks"]
            )
            mode_results["avg_time"] = (
                mode_results["total_time"] / mode_results["total_tasks"]
            )

        log.info(f"Result {mode_name}: SR={mode_results.get('success_rate', 0):.1%}")

        return mode_results

    def run_routed_mode(self):
        mode_results = {
            "success_count": 0,
            "total_tasks": 0,
            "total_time": 0,
            "token_usage": {
                "prompt_tokens": 0,
                "completion_tokens": 0,
                "total_tokens": 0,
            },
            "errors": [],
            "successful_indices": [],
            "routing_counts": {0: 0, 1: 0, 2: 0, 3: 0},
        }

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            futures = []
            for i, task in enumerate(self.tasks):
                emb = get_embedding(task["text"])
                cls = self.clf.predict([emb])[0]
                prompt_content = prompts[cls]
                mode_results["routing_counts"][cls] += 1
                futures.append(executor.submit(self.run_task, prompt_content, task, i))

            for future in concurrent.futures.as_completed(futures):
                res = future.result()
                mode_results["total_tasks"] += 1
                mode_results["total_time"] += res["time"]
                mode_results["token_usage"]["prompt_tokens"] += res["tokens"][
                    "prompt_tokens"
                ]
                mode_results["token_usage"]["completion_tokens"] += res["tokens"][
                    "completion_tokens"
                ]
                mode_results["token_usage"]["total_tokens"] += res["tokens"][
                    "total_tokens"
                ]

                log.info(f"Task {res['index'] + 1}/100 complete")
                if res["success"]:
                    mode_results["success_count"] += 1
                    mode_results["successful_indices"].append(res["index"])
                    log.info(f"Success ({res['time']:.2f}s)")
                else:
                    log.info(f"Error ({res['time']:.2f}s)")

                if "error" in res:
                    mode_results["errors"].append(res["error"])
                    log.error(f"Exception in task: {res['error']}")

        if mode_results["total_tasks"] > 0:
            mode_results["success_rate"] = (
                mode_results["success_count"] / mode_results["total_tasks"]
            )
            mode_results["avg_time"] = (
                mode_results["total_time"] / mode_results["total_tasks"]
            )

        log.info(
            f"Result routed: SR={mode_results.get('success_rate', 0):.1%}, Routing: {mode_results['routing_counts']}"
        )

        return mode_results

    def save_results(self, results, filename=None):
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"prompt_classifier_benchmark_results_{timestamp}.json"

        filepath = self.results_dir / filename

        serializable_results = {}
        for mode_name, data in results.items():
            serializable_results[mode_name] = {
                "success_rate": data.get("success_rate", 0),
                "success_count": data["success_count"],
                "total_tasks": data["total_tasks"],
                "avg_time": data.get("avg_time", 0),
                "total_time": data["total_time"],
                "token_usage": data["token_usage"],
                "errors": data["errors"],
                "successful_indices": data["successful_indices"],
            }
            if "routing_counts" in data:
                serializable_results[mode_name]["routing_counts"] = data[
                    "routing_counts"
                ]

        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(serializable_results, f, indent=2, ensure_ascii=False)

        log.info(f"Results saved to: {filepath}")
        return filepath


def main():
    model = "ibm-granite/granite-4.0-h-micro"

    benchmark = PromptClassifierBenchmark(model)

    log.info("Starting prompt classifier benchmark")
    results = benchmark.benchmark_modes()

    benchmark.save_results(results)


if __name__ == "__main__":
    main()
