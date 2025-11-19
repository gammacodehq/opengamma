import os
import sys

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import time
import json
import logging
import concurrent.futures
from datetime import datetime
from pathlib import Path
from datasets import load_dataset
from model import invoke_func

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
log = logging.getLogger(__name__)


system_prompt = f"""
Write a complete Python script using the python-pptx library to create a PowerPoint presentation (PPTX).
The script must:
- Use appropriate slide layouts (e.g., title slide, bullet slide).
- Save the presentation to a file named 'test.pptx'.
Output ONLY the Python code as plain text, without markdown, code block markers (e.g., ` + "```" + `python)

Example:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')"""

# Базовый промпт
basic_prompt = f"""
Generate Python code for python-pptx presentation. OUTPUT MUST BE PURE PYTHON CODE ONLY.

CRITICAL:
- NO comments (# comment)
- NO explanations
- NO markdown (```python or ```)
- NO text outside code
- NO example usage
- NO instructions

REQUIREMENTS:
- Use python-pptx
- Create multiple slides with different layouts  
- Save as 'test.pptx'
- Code must run without errors

NOW GENERATE PURE PYTHON CODE:"""

# Детальный промпт с примерами из документации
detailed_prompt = f"""
Generate Python code for python-pptx presentation. OUTPUT MUST BE PURE PYTHON CODE ONLY.

PYTHON-PPTX SPECIFICS:
- Use Presentation() to create presentation
- Use slide_layouts[] for slide types: 0=Title, 1=Title and Content, 5=Title only
- Use Inches() for measurements
- Add slides with prs.slides.add_slide(layout)
- Set text with shape.text = "text"
- Save with prs.save('test.pptx')

CODE STRUCTURE:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Title"
prs.save('test.pptx')

OUTPUT MUST BE PURE PYTHON CODE ONLY. NO COMMENTS. NO EXPLANATIONS."""

# Минималистичный промпт
minimal_prompt = f"""
Generate python-pptx code. Output ONLY Python code.

from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Title"
prs.save('test.pptx')

Continue with more slides:"""

# Строгий промпт
structured_prompt = f"""
PYTHON CODE ONLY. NO COMMENTS. FOLLOW TEMPLATE:

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "TITLE"
slide.placeholders[1].text = "SUBTITLE"

# ADD MORE SLIDES WITH DIFFERENT LAYOUTS

prs.save('test.pptx')"""


class PromptBenchmark:
    def __init__(self, model):
        self.model = model
        self.dataset = load_dataset("mikeoxmaul/opengamma-prs-dedup", streaming=True)
        self.results_dir = Path("results/prompt_benchmark")
        self.results_dir.mkdir(parents=True, exist_ok=True)

    def get_prompts_to_test(self):
        return {
            "original_prompt": system_prompt,
            "basic_prompt": basic_prompt,
            "detailed_prompt": detailed_prompt,
            "minimal_prompt": minimal_prompt,
            "structured_prompt": structured_prompt,
        }

    def run_task(self, prompt_content, task, i):
        id = 1
        try:
            start_time = time.time()
            result, token_stats = invoke_func(
                self.model, prompt_content, task["text"], id
            )
            execution_time = time.time() - start_time
            return {
                "success": result == 1,
                "time": execution_time,
                "tokens": token_stats,
                "index": i,
            }
        except Exception as e:
            return {
                "success": False,
                "time": 0,
                "tokens": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "error": str(e),
            }

    # Непосредственно бенчмарк
    def benchmark_prompts(self, prompts, num_tasks=10):
        results = {}

        for prompt_name, prompt_content in prompts.items():
            log.info(f"Testing prompt: {prompt_name}")

            prompt_results = {
                "success_count": 0,
                "total_tasks": 0,
                "total_time": 0,
                "token_usage": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "errors": [],
                "prompt_length": len(prompt_content),
                "successful_indices": [],
            }

            tasks = list(self.dataset["train"].take(num_tasks))

            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                futures = [
                    executor.submit(self.run_task, prompt_content, task, i)
                    for i, task in enumerate(tasks)
                ]
                for future in concurrent.futures.as_completed(futures):
                    res = future.result()
                    prompt_results["total_tasks"] += 1
                    prompt_results["total_time"] += res["time"]
                    prompt_results["token_usage"]["prompt_tokens"] += res["tokens"][
                        "prompt_tokens"
                    ]
                    prompt_results["token_usage"]["completion_tokens"] += res["tokens"][
                        "completion_tokens"
                    ]
                    prompt_results["token_usage"]["total_tokens"] += res["tokens"][
                        "total_tokens"
                    ]

                    log.info(f"Task {res["index"] + 1}/{num_tasks} complete")
                    if res["success"]:
                        prompt_results["success_count"] += 1
                        prompt_results["successful_indices"].append(res["index"])
                        log.info(f"Success ({res['time']:.2f}s)")
                    else:
                        log.info(f"Error ({res['time']:.2f}s)")

                    if "error" in res:
                        prompt_results["errors"].append(res["error"])
                        log.error(f"Exception in task: {res['error']}")

            if prompt_results["total_tasks"] > 0:
                prompt_results["success_rate"] = (
                    prompt_results["success_count"] / prompt_results["total_tasks"]
                )
                prompt_results["avg_time"] = (
                    prompt_results["total_time"] / prompt_results["total_tasks"]
                )

            results[prompt_name] = prompt_results
            log.info(
                f"Result {prompt_name}: SR={prompt_results.get('success_rate', 0):.1%}"
            )

        return results

    # Сохраняем результат в json в директорию results
    def save_results(self, results, filename=None):
        """Сохраняет результаты в JSON"""
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"prompt_benchmark_results_{timestamp}.json"

        filepath = self.results_dir / filename

        serializable_results = {}
        for prompt_name, data in results.items():
            serializable_results[prompt_name] = {
                "success_rate": data.get("success_rate", 0),
                "success_count": data["success_count"],
                "total_tasks": data["total_tasks"],
                "avg_time": data.get("avg_time", 0),
                "total_time": data["total_time"],
                "token_usage": data["token_usage"],
                "prompt_length": data["prompt_length"],
                "errors": data["errors"],
                "successful_indices": data["successful_indices"],
            }

        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(serializable_results, f, indent=2, ensure_ascii=False)

        log.info(f"Results are saved in : {filepath}")
        return filepath


def main():
    # Используем лучшую модель из предыдущего бенчмарка
    model = "ibm-granite/granite-4.0-h-micro"

    benchmark = PromptBenchmark(model)

    prompts = benchmark.get_prompts_to_test()

    log.info("Starting prompt benchmark")
    results = benchmark.benchmark_prompts(
        prompts, num_tasks=100
    )

    benchmark.save_results(results)


if __name__ == "__main__":
    main()
