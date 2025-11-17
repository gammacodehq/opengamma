import os
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from datasets import load_dataset
from datetime import datetime
import json
import logging
import concurrent.futures
from datetime import datetime
from pathlib import Path
from datasets import load_dataset
from model import invoke_func
from pathlib import Path
import time

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
log = logging.getLogger(__name__)

system_prompt = f"""
Write a complete Python script using the python-pptx library to create a PowerPoint presentation (PPTX).
The script must:
- Use appropriate slide layouts (e.g., title slide, bullet slide).
- Save the presentation to a file named 'test.pptx'.
Output ONLY the Python code as plain text, without markdown, code block markers (e.g., ` + "```" + `python)

Example:
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2

prs.save('test.pptx')"""


class ModelBenchmark:
    def __init__(self, system_prompt=None):
        self.system_prompt = system_prompt or self.get_default_prompt()
        self.dataset = load_dataset("mikeoxmaul/opengamma-prs", streaming=True)
        self.results_dir = Path("results/model_benchmark")
        self.results_dir.mkdir(parents=True, exist_ok=True)
        self.task_id = 0

    # Возвращием промпт для тестинга
    def get_default_prompt(self):
        return system_prompt

    def run_task(self, model_name, task, i):
        id = 1
        try:
            start_time = time.time()
            result, token_stats = invoke_func(
                model_name, self.system_prompt, task["text"], id
            )
            execution_time = time.time() - start_time
            return {
                "success": result == 1,
                "time": execution_time,
                "tokens": token_stats,
            }
        except Exception as e:
            return {
                "success": False,
                "time": 0,
                "tokens": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "error": str(e),
            }
        except Exception as e:
            return {
                "success": False,
                "time": 0,
                "tokens": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "error": str(e),
            }

    # Непосредственно бенчмарк
    def benchmark_models(self, models, num_tasks):
        results = {}

        for model_name in models:
            log.info(f"==== Testing model: {model_name}")

            model_results = {
                "success_count": 0,
                "total_tasks": 0,
                "total_time": 0,
                "token_usage": {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
                "errors": [],
            }

            tasks = list(self.dataset["train"].take(num_tasks))

            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                futures = [
                    executor.submit(self.run_task, model_name, task, i)
                    for i, task in enumerate(tasks)
                ]
                for future in concurrent.futures.as_completed(futures):
                    res = future.result()
                    model_results["total_tasks"] += 1
                    model_results["total_time"] += res["time"]
                    model_results["token_usage"]["prompt_tokens"] += res["tokens"][
                        "prompt_tokens"
                    ]
                    model_results["token_usage"]["completion_tokens"] += res["tokens"][
                        "completion_tokens"
                    ]
                    model_results["token_usage"]["total_tokens"] += res["tokens"][
                        "total_tokens"
                    ]
                    if res["success"]:
                        model_results["success_count"] += 1
                    if "error" in res:
                        model_results["errors"].append(res["error"])

            # Считаем метрки
            if model_results["total_tasks"] > 0:
                model_results["success_rate"] = (
                    model_results["success_count"] / model_results["total_tasks"]
                )
                model_results["avg_time"] = (
                    model_results["total_time"] / model_results["total_tasks"]
                )

            results[model_name] = model_results
            log.info(
                f"==== Result {model_name}: SR={model_results.get('success_rate', 0):.1%}"
            )

        return results

    # Сохраняем рещультат в json в директорию results
    def save_results(self, results, filename=None):
        """Сохраняет результаты в JSON"""
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"benchmark_results_{timestamp}.json"

        filepath = self.results_dir / filename

        # Конвертируем в подходящий для json.dump формат
        serializable_results = {}
        for model, data in results.items():
            serializable_results[model] = {
                "success_rate": data.get("success_rate", 0),
                "success_count": data["success_count"],
                "total_tasks": data["total_tasks"],
                "avg_time": data.get("avg_time", 0),
                "total_time": data["total_time"],
                "token_usage": data["token_usage"],
                "errors": data["errors"],
            }

        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(serializable_results, f, indent=2, ensure_ascii=False)

        log.info(f"==== Results are saved in : {filepath}")
        return filepath

def main():
    models_to_test = [
        "openai/gpt-oss-20b",  # Current
        "x-ai/grok-code-fast-1",  # #1 in Programming on openrouter
        "ibm-granite/granite-4.0-h-micro",  # Cheaper and newer than gpt-oss-20b
    ]

    benchmark = ModelBenchmark()

    log.info("==== Starting benchmark")
    num_tasks = 100
    results = benchmark.benchmark_models(models_to_test, num_tasks)
    benchmark.save_results(results)


if __name__ == "__main__":
    main()
